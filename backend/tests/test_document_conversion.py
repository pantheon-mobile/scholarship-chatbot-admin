from io import BytesIO

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

from app.services.document_conversion import convert_docx, convert_pptx, convert_xlsx


def test_xlsx_is_split_into_visible_markdown_sheets():
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "入力"
    sheet.append(["項目", "金額"])
    sheet.append(["授業料", 1000])
    hidden = workbook.create_sheet("非表示")
    hidden.sheet_state = "hidden"
    hidden.append(["秘密"])
    output = BytesIO()
    workbook.save(output)

    documents = convert_xlsx(output.getvalue(), "sample.xlsx")

    assert len(documents) == 1
    assert documents[0].name == "sheet-001.md"
    assert "## シート: 入力" in documents[0].markdown
    assert "| 授業料 | 1000 |" in documents[0].markdown


def test_docx_keeps_headings_paragraphs_and_tables():
    document = Document()
    document.add_heading("奨学金案内", level=1)
    document.add_paragraph("申請してください。")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "区分"
    table.cell(0, 1).text = "金額"
    table.cell(1, 0).text = "第一種"
    table.cell(1, 1).text = "50000"
    output = BytesIO()
    document.save(output)

    converted = convert_docx(output.getvalue(), "guide.docx")[0].markdown

    assert "## 奨学金案内" in converted
    assert "申請してください。" in converted
    assert "| 第一種 | 50000 |" in converted


def test_pptx_keeps_slide_order_text_and_tables():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "制度概要"
    textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(5), Inches(1))
    textbox.text = "対象者は学生です。"
    output = BytesIO()
    presentation.save(output)

    converted = convert_pptx(output.getvalue(), "guide.pptx")[0].markdown

    assert "## 1. 制度概要" in converted
    assert "対象者は学生です。" in converted
