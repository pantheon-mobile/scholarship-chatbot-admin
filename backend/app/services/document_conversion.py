from __future__ import annotations

import base64
import csv
import io
import json
import os
import re
from dataclasses import dataclass
from urllib.parse import urljoin, urlsplit, urlunsplit
from urllib.robotparser import RobotFileParser

import boto3
import fitz
import requests
from bs4 import BeautifulSoup
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation


@dataclass(frozen=True)
class ConvertedDocument:
    name: str
    markdown: str
    source_url: str | None = None
    metadata: dict | None = None


def _table_markdown(rows: list[list[str]]) -> str:
    if not rows:
        return ""
    width = max(len(row) for row in rows)
    normalized = [row + [""] * (width - len(row)) for row in rows]
    escaped = [[str(value or "").replace("|", "\\|").replace("\n", "<br>") for value in row]
               for row in normalized]
    return "\n".join([
        "| " + " | ".join(escaped[0]) + " |",
        "| " + " | ".join(["---"] * width) + " |",
        *("| " + " | ".join(row) + " |" for row in escaped[1:]),
    ])


def convert_xlsx(content: bytes, name: str) -> list[ConvertedDocument]:
    workbook = load_workbook(io.BytesIO(content), data_only=True, read_only=True)
    documents = []
    for index, sheet in enumerate(workbook.worksheets, start=1):
        if sheet.sheet_state != "visible":
            continue
        rows = [["" if value is None else str(value) for value in row]
                for row in sheet.iter_rows(values_only=True)]
        while rows and not any(cell for cell in rows[-1]):
            rows.pop()
        if not rows:
            continue
        documents.append(ConvertedDocument(
            name=f"sheet-{index:03d}.md",
            markdown=f"# {name}\n\n## シート: {sheet.title}\n\n{_table_markdown(rows)}",
            metadata={"sheet_name": sheet.title, "sheet_index": index},
        ))
    workbook.close()
    return documents


def convert_docx(content: bytes, name: str) -> list[ConvertedDocument]:
    document = Document(io.BytesIO(content))
    blocks: list[str] = [f"# {name}"]
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style = (paragraph.style.name or "").lower() if paragraph.style else ""
        if style.startswith("heading"):
            level = int("".join(character for character in style if character.isdigit()) or "1")
            blocks.append(f"{'#' * min(level + 1, 6)} {text}")
        elif "bullet" in style or "箇条書き" in style:
            blocks.append(f"- {text}")
        else:
            blocks.append(text)
    for table in document.tables:
        blocks.append(_table_markdown([[cell.text.strip() for cell in row.cells] for row in table.rows]))
    return [ConvertedDocument("source.md", "\n\n".join(blocks))]


def convert_pptx(content: bytes, name: str) -> list[ConvertedDocument]:
    presentation = Presentation(io.BytesIO(content))
    blocks = [f"# {name}"]
    for slide_number, slide in enumerate(presentation.slides, start=1):
        title = slide.shapes.title.text.strip() if slide.shapes.title else f"スライド {slide_number}"
        blocks.extend([f"## {slide_number}. {title}", f"<!-- slide {slide_number} -->"])
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                blocks.append(_table_markdown(
                    [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                ))
            elif getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text and text != title:
                    blocks.append(text)
        try:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            notes = ""
        if notes:
            blocks.extend(["### 発表者ノート", notes])
    return [ConvertedDocument("source.md", "\n\n".join(blocks))]


def _vision_page_markdown(page, page_number: int) -> str:
    model_id = os.getenv("PDF_VISION_MODEL_ID", "").strip()
    if not model_id:
        raise RuntimeError(
            "画像中心のPDFです。PDF_VISION_MODEL_IDを設定してVision変換を有効にしてください。"
        )
    pixmap = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False)
    image = base64.b64encode(pixmap.tobytes("png")).decode("ascii")
    client = boto3.client("bedrock-runtime", region_name=os.getenv("AWS_REGION", "ap-northeast-1"))
    response = client.invoke_model(
        modelId=model_id,
        contentType="application/json",
        accept="application/json",
        body=json.dumps({
            "anthropic_version": "bedrock-2023-05-31",
            "max_tokens": 4000,
            "temperature": 0,
            "messages": [{"role": "user", "content": [
                {"type": "image", "source": {"type": "base64", "media_type": "image/png", "data": image}},
                {"type": "text", "text": "このページを、見出し・箇条書き・表を保った日本語Markdownに変換してください。Markdown本文だけを返してください。"},
            ]}],
        }).encode("utf-8"),
    )
    body = json.loads(response["body"].read())
    return "\n".join(item["text"] for item in body.get("content", []) if item.get("type") == "text")


def convert_pdf(content: bytes, name: str) -> list[ConvertedDocument]:
    pdf = fitz.open(stream=content, filetype="pdf")
    page_texts = [page.get_text("text").strip() for page in pdf]
    image_count = sum(len(page.get_images(full=True)) for page in pdf)
    extracted_characters = sum(len(text) for text in page_texts)
    chars_per_page = extracted_characters / max(len(pdf), 1)
    use_vision = image_count > 0 and chars_per_page < float(os.getenv("PDF_VISION_MIN_CHARS_PER_PAGE", "100"))
    method = "VISION_MARKDOWN" if use_vision else "TEXT_MARKDOWN"
    pages = []
    for index, page in enumerate(pdf, start=1):
        text = _vision_page_markdown(page, index) if use_vision else page_texts[index - 1]
        pages.append(f"## ページ {index}\n\n{text}")
    metadata = {
        "conversion_method": method,
        "selection_reason": "image_heavy_low_text_density" if use_vision else "text_extraction_sufficient",
        "page_count": len(pdf),
        "image_count": image_count,
        "extracted_characters": extracted_characters,
        "characters_per_page": round(chars_per_page, 1),
    }
    return [ConvertedDocument("source.md", f"# {name}\n\n" + "\n\n".join(pages), metadata=metadata)]


def convert_plain_text(content: bytes, name: str) -> list[ConvertedDocument]:
    text = content.decode("utf-8-sig", errors="replace")
    if name.lower().endswith(".csv"):
        rows = list(csv.reader(io.StringIO(text)))
        text = _table_markdown(rows)
    return [ConvertedDocument("source.md", f"# {name}\n\n{text}")]


def _normalize_url(url: str, base: str = "") -> str:
    parts = urlsplit(urljoin(base, url))
    if parts.scheme not in {"http", "https"} or not parts.hostname:
        return ""
    return urlunsplit((parts.scheme.lower(), parts.netloc.lower(), parts.path or "/", parts.query, ""))


def crawl_website(root_url: str) -> list[ConvertedDocument]:
    root = _normalize_url(root_url)
    root_parts = urlsplit(root)
    root_path = root_parts.path if root_parts.path.endswith("/") else root_parts.path.rsplit("/", 1)[0] + "/"
    max_pages = int(os.getenv("WEB_CRAWL_MAX_PAGES", "500"))
    max_depth = int(os.getenv("WEB_CRAWL_MAX_DEPTH", "5"))
    timeout = float(os.getenv("WEB_CRAWL_TIMEOUT_SECONDS", "20"))
    session = requests.Session()
    session.headers["User-Agent"] = os.getenv("WEB_CRAWL_USER_AGENT", "ScholarshipChatbotCrawler/1.0")
    robots = RobotFileParser(f"{root_parts.scheme}://{root_parts.netloc}/robots.txt")
    try:
        robots.read()
    except Exception:
        robots = None
    pending = [(root, 0)]
    visited: set[str] = set()
    documents: list[ConvertedDocument] = []
    while pending and len(documents) < max_pages:
        url, depth = pending.pop(0)
        if url in visited or depth > max_depth:
            continue
        visited.add(url)
        parts = urlsplit(url)
        if parts.hostname != root_parts.hostname or not parts.path.startswith(root_path):
            continue
        if robots and not robots.can_fetch(session.headers["User-Agent"], url):
            continue
        response = session.get(url, timeout=timeout)
        response.raise_for_status()
        if "html" not in response.headers.get("Content-Type", "").lower():
            continue
        soup = BeautifulSoup(response.text, "lxml")
        for node in soup.select("script,style,nav,header,footer,aside,form,iframe"):
            node.decompose()
        main = soup.find("main") or soup.find("article") or soup.body
        if main is None:
            continue
        title = soup.title.get_text(" ", strip=True) if soup.title else url
        text = re.sub(r"\n{3,}", "\n\n", main.get_text("\n", strip=True))
        if len(text) >= 20:
            page_id = __import__("hashlib").sha256(url.encode()).hexdigest()[:20]
            documents.append(ConvertedDocument(
                f"web-{page_id}.md", f"# {title}\n\n元URL: {url}\n\n{text}",
                source_url=url, metadata={"crawl_depth": depth},
            ))
        for anchor in main.find_all("a", href=True):
            child = _normalize_url(anchor["href"], url)
            if child and child not in visited:
                pending.append((child, depth + 1))
    if not documents:
        raise RuntimeError("Webサイトから本文を取得できませんでした。")
    return documents
